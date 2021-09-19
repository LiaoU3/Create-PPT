# 2021/09/13 version 1.2
# Author : Vincent

from pptx import Presentation
import datetime
import os

# get the current n'th week of the month
def get_week_of_month():
    today = datetime.datetime.now()
    begin = int(datetime.date(today.year, today.month, 1).strftime("%W"))
    end = int(datetime.date(today.year, today.month, today.day).strftime("%W"))

    return end - begin + 1

def main():
    # set your name here
    name = 'Vincent'

    # you will get the correct weekday, if you create file on Monday(which is the next monday)
    if datetime.datetime.today().weekday() == 0 :
        week = get_week_of_month() - 1
    else:
        week = get_week_of_month()

    # set the correct file name according to the current date
    file_name = str(datetime.datetime.now().year) + '-' + str(datetime.datetime.now().month) + 'M-' +str(week) + 'W-週報'

    # check if the file is already exist
    keep_going = 'y'
    if os.path.isfile(file_name + '.pptx'):
        keep_going = input('The file is already exist, do you want to delete it and create a new file ? (y/n) ')
    if keep_going == 'y':
        # Create pptx object
        prs = Presentation('src/Template.pptx')
        # Set the title of first page 
        title = prs.slides[0].shapes.title
        title.text = file_name
        # Set the subtitle of first page 
        subtitle = prs.slides[0].placeholders[1]
        subtitle.text = name
        prs.save(file_name + '.pptx')
        print('Sate : File Created -> ' + file_name + '.pptx')
    else:
        print("State : Didn't do anything")

    input("Press any key to exit...")

if __name__ == '__main__':
    main()